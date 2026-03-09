"""
RAG Web App — Assistant documentaire avec Mistral.

Interface en 3 onglets :
  📂 Sources   — Upload, prévisualisation et gestion des fichiers
  ⚙️ Pipeline  — Lancement et suivi du pipeline RAG en arrière-plan
  💬 Chat      — Assistant RAG conversationnel

Usage :
    streamlit run streamlit/app.py
"""

from __future__ import annotations

import sys
import time
import shutil
import logging
import base64
from pathlib import Path

import streamlit as st

# ── Chemins ──────────────────────────────────────────────────────────────────
APP_DIR = Path(__file__).parent
ROOT_DIR = APP_DIR.parent
sys.path.insert(0, str(ROOT_DIR))

from config import (
    INPUTS_DIR, MARKDOWN_DIR, FAISS_INDEX_FILE, FAISS_METADATA_FILE,
    AVAILABLE_MODELS, DEFAULT_MODEL, GENERATION_PARAMS,
    SUPPORTED_EXTENSIONS, DOCLING_EXTENSIONS, AUDIO_EXTENSIONS,
    MISTRAL_API_KEY,
)

# ── Configuration Streamlit ───────────────────────────────────────────────────
st.set_page_config(
    page_title="RAG Assistant — Mistral",
    page_icon="🧠",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ── CSS dynamique (dark / light) ─────────────────────────────────────────────

def _inject_theme_css(dark: bool) -> None:
    """Injecte le CSS de thème selon le mode choisi."""
    if dark:
        css = """
<style>
:root {
  --bg-dark:     #0d0f14;
  --bg-card:     #161b27;
  --bg-card2:    #1c2232;
  --bg-card3:    #1a1f2f;
  --accent:      #6c63ff;
  --accent2:     #00d4aa;
  --accent3:     #ff6b6b;
  --text-main:   #e8eaf0;
  --text-muted:  #8b92a5;
  --border:      rgba(108,99,255,0.2);
  --radius:      14px;
  --shadow:      0 4px 24px rgba(0,0,0,0.4);
  --log-bg:      #0a0c12;
  --log-border:  #2d3148;
  --log-text:    #b8c0d8;
  --upload-bg:   rgba(108,99,255,0.06);
  --upload-border: rgba(108,99,255,0.25);
  --upload-hover: rgba(108,99,255,0.12);
  --glass-bg:    rgba(22,27,39,0.7);
  --glass-border: rgba(108,99,255,0.15);
}
.stApp {
  background: linear-gradient(135deg, #0d0f14 0%, #111827 60%, #0d1117 100%) !important;
  color: var(--text-main);
}
section[data-testid="stSidebar"] { background: var(--bg-card) !important; border-right: 1px solid var(--border); }
section[data-testid="stSidebar"] * { color: var(--text-main) !important; }
"""
    else:
        css = """
<style>
:root {
  --bg-dark:     #f5f7fb;
  --bg-card:     #ffffff;
  --bg-card2:    #f0f2f8;
  --bg-card3:    #e8ecf4;
  --accent:      #5b52e8;
  --accent2:     #00a98f;
  --accent3:     #e05555;
  --text-main:   #1a1d2e;
  --text-muted:  #5a6275;
  --border:      rgba(91,82,232,0.18);
  --radius:      14px;
  --shadow:      0 4px 24px rgba(0,0,0,0.08);
  --log-bg:      #f0f2f8;
  --log-border:  #d0d5e8;
  --log-text:    #3a4060;
  --upload-bg:   rgba(91,82,232,0.04);
  --upload-border: rgba(91,82,232,0.2);
  --upload-hover: rgba(91,82,232,0.08);
  --glass-bg:    rgba(255,255,255,0.7);
  --glass-border: rgba(91,82,232,0.12);
}
.stApp {
  background: linear-gradient(135deg, #f5f7fb 0%, #eef1f9 60%, #f0f4fc 100%) !important;
  color: var(--text-main);
}
section[data-testid="stSidebar"] { background: var(--bg-card) !important; border-right: 1px solid var(--border); }
section[data-testid="stSidebar"] * { color: var(--text-main) !important; }
.hero-upload { background: var(--bg-card) !important; }
.rag-card { background: var(--bg-card) !important; }
.file-item { background: var(--bg-card2) !important; }
.step-row  { background: var(--bg-card2) !important; }
.pipeline-card { background: var(--bg-card) !important; }
.chat-welcome { background: var(--bg-card) !important; }
"""

    # Partie commune (indépendante du thème)
    common = """
/* ── TABS ────────────────────────────────────────────────── */
.stTabs [data-baseweb="tab-list"] {
  background: var(--bg-card);
  border-radius: 16px;
  padding: 6px;
  gap: 6px;
  border: 1px solid var(--border);
  box-shadow: var(--shadow);
}
.stTabs [data-baseweb="tab"] {
  background: transparent;
  border-radius: 12px;
  color: var(--text-muted);
  font-weight: 600;
  font-size: 0.95rem;
  padding: 0.6rem 1.4rem;
  transition: all 0.25s ease;
}
.stTabs [aria-selected="true"] {
  background: linear-gradient(135deg, var(--accent), #9b59f5) !important;
  color: white !important;
  box-shadow: 0 4px 15px rgba(108,99,255,0.3);
}

/* ── CARDS ────────────────────────────────────────────────── */
.rag-card {
  border: 1px solid var(--border);
  border-radius: var(--radius);
  padding: 1.4rem 1.6rem;
  margin-bottom: 1.2rem;
  box-shadow: var(--shadow);
  backdrop-filter: blur(8px);
}

/* ── UPLOAD HERO ─────────────────────────────────────────── */
.hero-upload {
  border: 2px dashed var(--upload-border);
  border-radius: 20px;
  padding: 2.5rem 2rem;
  margin-bottom: 2rem;
  text-align: center;
  background: var(--upload-bg);
  transition: all 0.3s ease;
  box-shadow: var(--shadow);
}
.hero-upload:hover {
  background: var(--upload-hover);
  border-color: var(--accent);
  box-shadow: 0 6px 30px rgba(108,99,255,0.2);
}
.hero-upload h1 {
  font-size: 2rem;
  font-weight: 800;
  background: linear-gradient(90deg, var(--accent), var(--accent2));
  -webkit-background-clip: text;
  -webkit-text-fill-color: transparent;
  margin: 0 0 0.3rem;
}
.hero-upload .hero-subtitle {
  color: var(--text-muted);
  font-size: 1rem;
  margin: 0 0 0.8rem;
}
.hero-upload .hero-icon {
  font-size: 3rem;
  margin-bottom: 0.8rem;
  filter: drop-shadow(0 2px 8px rgba(108,99,255,0.3));
}
.hero-upload .supported-types {
  display: flex;
  flex-wrap: wrap;
  gap: 6px;
  justify-content: center;
  margin-top: 1rem;
}
.hero-upload .type-tag {
  background: rgba(108,99,255,0.1);
  color: var(--accent);
  padding: 3px 10px;
  border-radius: 20px;
  font-size: 0.72rem;
  font-weight: 600;
  letter-spacing: 0.04em;
  border: 1px solid rgba(108,99,255,0.15);
}

/* ── PIPELINE CARD ───────────────────────────────────────── */
.pipeline-card {
  border: 1px solid var(--border);
  border-radius: 20px;
  padding: 2rem 2.2rem;
  margin-bottom: 1.5rem;
  box-shadow: var(--shadow);
  backdrop-filter: blur(8px);
}
.pipeline-card h1 {
  font-size: 2rem;
  font-weight: 800;
  background: linear-gradient(90deg, var(--accent), var(--accent2));
  -webkit-background-clip: text;
  -webkit-text-fill-color: transparent;
  margin: 0 0 0.4rem;
  text-align: center;
}
.pipeline-card .pipeline-subtitle {
  color: var(--text-muted);
  font-size: 1rem;
  text-align: center;
  margin-bottom: 1.5rem;
}

/* ── BADGES ──────────────────────────────────────────────── */
.badge { display:inline-block; padding:.25rem .7rem; border-radius:20px; font-size:.78rem; font-weight:700; letter-spacing:.05em; text-transform:uppercase; }
.badge-green  { background:rgba(0,212,170,.15); color:#00a98f; border:1px solid rgba(0,212,170,.3);}
.badge-yellow { background:rgba(255,190,50,.15); color:#b8860b; border:1px solid rgba(255,190,50,.3);}
.badge-red    { background:rgba(255,107,107,.15); color:var(--accent3); border:1px solid rgba(255,107,107,.3);}
.badge-blue   { background:rgba(108,99,255,.15); color:var(--accent); border:1px solid rgba(108,99,255,.3);}

/* ── FILES ───────────────────────────────────────────────── */
.file-item {
  display: flex;
  align-items: center;
  gap: .7rem;
  border: 1px solid var(--border);
  border-radius: 12px;
  padding: .65rem 1rem;
  margin-bottom: .5rem;
  transition: all 0.25s ease;
}
.file-item:hover { border-color: var(--accent); transform: translateX(4px); }
.file-icon { font-size: 1.4rem; }
.file-name { font-weight: 600; flex: 1; color: var(--text-main); }
.file-size { color: var(--text-muted); font-size: .82rem; }
.file-type-tag {
  background: rgba(108,99,255,0.1);
  color: var(--accent);
  padding: 2px 8px;
  border-radius: 12px;
  font-size: 0.7rem;
  font-weight: 600;
  text-transform: uppercase;
}

/* empty state */
.empty-state {
  text-align: center;
  padding: 3rem 2rem;
  border: 1px solid var(--border);
  border-radius: 16px;
}
.empty-state .empty-icon { font-size: 3.5rem; margin-bottom: 1rem; opacity: 0.7; }
.empty-state .empty-text { color: var(--text-muted); font-size: 1rem; }

/* ── STEPS ───────────────────────────────────────────────── */
.step-row {
  display: flex;
  align-items: center;
  gap: 1rem;
  padding: .7rem 1.2rem;
  border-radius: 12px;
  margin-bottom: .5rem;
  border: 1px solid var(--border);
  transition: all 0.3s ease;
}
.step-row.active  { border-color: var(--accent); background: rgba(108,99,255,.08); box-shadow: 0 0 12px rgba(108,99,255,0.1); }
.step-row.done    { border-color: #00a98f; background: rgba(0,169,143,.06); }
.step-row.pending { opacity: 0.45; }
.step-num {
  width: 32px; height: 32px; border-radius: 50%;
  display: flex; align-items: center; justify-content: center;
  font-weight: 700; font-size: .85rem;
  background: var(--bg-card);
  border: 2px solid var(--border);
  flex-shrink: 0;
  transition: all 0.3s ease;
}
.step-row.active .step-num { background: var(--accent); border-color: var(--accent); color: white; box-shadow: 0 0 10px rgba(108,99,255,0.4); }
.step-row.done   .step-num { background: #00a98f; border-color: #00a98f; color: #fff; }
.step-label { font-weight: 600; flex: 1; }
.step-desc  { color: var(--text-muted); font-size: .82rem; }

/* ── BUTTONS ─────────────────────────────────────────────── */
div.stButton > button { font-weight: 600; border-radius: 10px; transition: all .2s; }
div.stButton > button[kind="secondary"] { background-color: var(--bg-card) !important; color: var(--text-main) !important; border: 1px solid var(--border) !important; }
div.stButton > button[kind="secondary"]:hover { border-color: var(--accent) !important; color: var(--accent) !important; }
div.stButton > button[kind="primary"] { background: linear-gradient(135deg, var(--accent), #9b59f5) !important; color: white !important; border: none !important; }
div.stButton > button[kind="primary"]:hover { box-shadow: 0 4px 15px rgba(108,99,255,0.4) !important; transform: translateY(-1px); }

/* Big launch button */
.launch-btn-wrap div.stButton > button[kind="primary"] {
  font-size: 1.15rem !important;
  padding: 0.85rem 2rem !important;
  border-radius: 14px !important;
  font-weight: 700 !important;
  letter-spacing: 0.02em;
  background: linear-gradient(135deg, var(--accent), #9b59f5, var(--accent2)) !important;
  background-size: 200% 200% !important;
  animation: gradientShift 3s ease infinite;
  box-shadow: 0 6px 24px rgba(108,99,255,0.35) !important;
}
.launch-btn-wrap div.stButton > button[kind="primary"]:hover {
  box-shadow: 0 8px 30px rgba(108,99,255,0.5) !important;
  transform: translateY(-2px);
}
@keyframes gradientShift {
  0% { background-position: 0% 50%; }
  50% { background-position: 100% 50%; }
  100% { background-position: 0% 50%; }
}

/* ── CHAT ────────────────────────────────────────────────── */
div[data-testid="stChatMessage"] {
  background: var(--bg-card) !important;
  border: 1px solid var(--border) !important;
  border-radius: 14px !important;
  margin-bottom: .7rem !important;
  box-shadow: 0 2px 12px rgba(0,0,0,0.1) !important;
}
.chat-welcome {
  border: 1px solid var(--border);
  border-radius: 20px;
  padding: 2.5rem;
  text-align: center;
  margin-bottom: 1.5rem;
  box-shadow: var(--shadow);
}
.chat-welcome h1 {
  font-size: 2rem;
  font-weight: 800;
  background: linear-gradient(90deg, var(--accent), var(--accent2));
  -webkit-background-clip: text;
  -webkit-text-fill-color: transparent;
  margin: 0 0 0.4rem;
}
.chat-welcome .chat-subtitle {
  color: var(--text-muted);
  font-size: 1rem;
}
.chat-feature-grid {
  display: grid;
  grid-template-columns: repeat(3, 1fr);
  gap: 1rem;
  margin-top: 1.5rem;
}
.chat-feature {
  border: 1px solid var(--border);
  border-radius: 14px;
  padding: 1.2rem;
  text-align: center;
  transition: all 0.25s ease;
}
.chat-feature:hover { border-color: var(--accent); transform: translateY(-2px); }
.chat-feature .feat-icon { font-size: 1.8rem; margin-bottom: 0.5rem; }
.chat-feature .feat-title { font-weight: 700; font-size: 0.9rem; color: var(--text-main); }
.chat-feature .feat-desc { color: var(--text-muted); font-size: 0.78rem; margin-top: 0.3rem; }

/* ── LOGS ────────────────────────────────────────────────── */
.log-box {
  background: var(--log-bg);
  border: 1px solid var(--log-border);
  border-radius: 12px;
  padding: 1.2rem;
  font-family: 'JetBrains Mono', 'Fira Code', monospace;
  font-size: .78rem;
  color: var(--log-text);
  max-height: 320px;
  overflow-y: auto;
  white-space: pre-wrap;
  line-height: 1.7;
}

/* ── PREVIEW ─────────────────────────────────────────────── */
.preview-box {
  border: 1px solid var(--border);
  border-radius: 12px;
  overflow: hidden;
  background: var(--bg-card2);
}

/* ── SECTION HEADERS ─────────────────────────────────────── */
.section-header {
  display: flex;
  align-items: center;
  gap: 0.6rem;
  margin-bottom: 1rem;
}
.section-header h3 {
  margin: 0;
  font-weight: 700;
  font-size: 1.15rem;
}
.section-header .section-count {
  background: rgba(108,99,255,0.12);
  color: var(--accent);
  padding: 2px 10px;
  border-radius: 20px;
  font-size: 0.78rem;
  font-weight: 700;
}

/* ── SCROLLBAR ───────────────────────────────────────────── */
::-webkit-scrollbar { width: 5px; height: 5px; }
::-webkit-scrollbar-track { background: transparent; }
::-webkit-scrollbar-thumb { background: #a0a8c0; border-radius: 4px; }

/* ── DIVIDER ─────────────────────────────────────────────── */
.custom-divider {
  border: none;
  border-top: 1px solid var(--border);
  margin: 1.5rem 0;
}

/* Stats pills */
.stat-pill {
  display: inline-flex;
  align-items: center;
  gap: 6px;
  background: var(--bg-card2);
  border: 1px solid var(--border);
  border-radius: 20px;
  padding: 4px 12px;
  font-size: 0.82rem;
  font-weight: 600;
  margin-right: 8px;
}

/* ── VALIDATED FILE STATE ────────────────────────────────── */
.file-item.validated {
  border-color: #00a98f !important;
  background: rgba(0,169,143,0.06) !important;
}
.file-item.validated .file-name::after {
  content: ' ✅';
  font-size: 0.85rem;
}
.validated-badge {
  display: inline-flex; align-items: center; gap: 6px;
  background: rgba(0,212,170,0.12); color: #00a98f;
  border: 1px solid rgba(0,212,170,0.3);
  border-radius: 20px; padding: 4px 14px;
  font-size: 0.82rem; font-weight: 700;
  letter-spacing: 0.03em;
}

/* ── STEPS GRID (Pipeline) ──────────────────────────────── */
.steps-grid {
  display: grid;
  grid-template-columns: repeat(4, 1fr);
  gap: 1rem;
  margin-bottom: 1.5rem;
}
.step-card {
  border: 1px solid var(--border);
  border-radius: 14px;
  padding: 1.2rem 1rem;
  text-align: center;
  transition: all 0.3s ease;
  position: relative;
  overflow: hidden;
}
.step-card.active {
  border-color: var(--accent);
  background: rgba(108,99,255,0.08);
  box-shadow: 0 0 16px rgba(108,99,255,0.15);
}
.step-card.done {
  border-color: #00a98f;
  background: rgba(0,169,143,0.06);
}
.step-card.pending { opacity: 0.45; }
.step-card .step-icon { font-size: 1.8rem; margin-bottom: 0.5rem; }
.step-card .step-num-badge {
  width: 28px; height: 28px; border-radius: 50%;
  display: inline-flex; align-items: center; justify-content: center;
  font-weight: 700; font-size: 0.8rem;
  background: var(--bg-card); border: 2px solid var(--border);
  margin-bottom: 0.6rem;
}
.step-card.active .step-num-badge { background: var(--accent); border-color: var(--accent); color: white; }
.step-card.done .step-num-badge { background: #00a98f; border-color: #00a98f; color: white; }
.step-card .step-card-label { font-weight: 700; font-size: 0.95rem; margin-bottom: 0.3rem; }
.step-card .step-card-desc { color: var(--text-muted); font-size: 0.78rem; }

/* ── CHAT ALIGNMENT ─────────────────────────────────────── */
div[data-testid="stChatMessage"][aria-label="user"] {
  margin-left: 20% !important;
  margin-right: 0 !important;
  background: linear-gradient(135deg, var(--accent), #9b59f5) !important;
  color: white !important;
}
div[data-testid="stChatMessage"][aria-label="user"] * {
  color: white !important;
}
div[data-testid="stChatMessage"][aria-label="assistant"] {
  margin-right: 20% !important;
  margin-left: 0 !important;
}

/* ── SUCCESS POPUP ──────────────────────────────────────── */
.success-popup-overlay {
  position: fixed; top: 0; left: 0; width: 100vw; height: 100vh;
  background: rgba(0,0,0,0.5); backdrop-filter: blur(6px);
  z-index: 9999; display: flex; align-items: center; justify-content: center;
}
.success-popup {
  background: var(--bg-card);
  border: 2px solid #00a98f;
  border-radius: 24px;
  padding: 3rem 3.5rem;
  text-align: center;
  box-shadow: 0 20px 60px rgba(0,0,0,0.4);
  animation: popupIn 0.5s cubic-bezier(0.34, 1.56, 0.64, 1);
  max-width: 500px;
}
@keyframes popupIn {
  0% { transform: scale(0.5); opacity: 0; }
  100% { transform: scale(1); opacity: 1; }
}
.success-popup .popup-icon { font-size: 4rem; margin-bottom: 1rem; }
.success-popup h2 {
  font-size: 1.6rem; font-weight: 800; margin: 0 0 0.5rem;
  background: linear-gradient(90deg, #00a98f, var(--accent));
  -webkit-background-clip: text; -webkit-text-fill-color: transparent;
}
.success-popup p { color: var(--text-muted); margin-bottom: 1.5rem; }

/* ── PULSE BUTTON ANIMATION ─────────────────────────────── */
@keyframes pulseBtn {
  0% { transform: scale(1); box-shadow: 0 6px 24px rgba(108,99,255,0.35); }
  50% { transform: scale(1.04); box-shadow: 0 8px 32px rgba(108,99,255,0.5); }
  100% { transform: scale(1); box-shadow: 0 6px 24px rgba(108,99,255,0.35); }
}
.pulse-btn div.stButton > button[kind="primary"] {
  animation: pulseBtn 1.8s ease-in-out infinite, gradientShift 3s ease infinite !important;
}

/* ── VISUALIZATION ──────────────────────────────────────── */
.viz-card {
  border: 1px solid var(--border);
  border-radius: 16px;
  padding: 1.5rem;
  margin-bottom: 1rem;
  box-shadow: var(--shadow);
  transition: all 0.3s ease;
}
.viz-card:hover { border-color: var(--accent); transform: translateY(-2px); }
.viz-card h3 { margin: 0 0 0.8rem; font-weight: 700; font-size: 1.05rem; }
.viz-stat { display: flex; justify-content: space-between; align-items: center; padding: 0.5rem 0; border-bottom: 1px solid var(--border); }
.viz-stat:last-child { border-bottom: none; }
.viz-stat-label { color: var(--text-muted); font-size: 0.88rem; }
.viz-stat-value { font-weight: 700; font-size: 0.95rem; color: var(--accent); }
</style>
"""
    st.markdown(css + common, unsafe_allow_html=True)

# ── Logging ───────────────────────────────────────────────────────────────────
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


# ── Helpers ───────────────────────────────────────────────────────────────────

def human_size(n_bytes: int) -> str:
    for unit in ("o", "Ko", "Mo", "Go"):
        if n_bytes < 1024:
            return f"{n_bytes:.0f} {unit}"
        n_bytes /= 1024
    return f"{n_bytes:.1f} To"


def file_icon(suffix: str) -> str:
    s = suffix.lower()
    if s == ".pdf":       return "📕"
    if s == ".docx":      return "📘"
    if s == ".pptx":      return "📊"
    if s == ".xlsx":      return "📗"
    if s in {".html",".htm"}: return "🌐"
    if s in {".png",".jpg",".jpeg",".webp",".tiff",".tif",".bmp"}: return "🖼️"
    if s in {".wav",".mp3",".m4a",".ogg",".flac",".webm"}: return "🎵"
    if s == ".tex":       return "📐"
    if s == ".vtt":       return "💬"
    if s == ".csv":       return "📋"
    return "📄"


def index_ready() -> bool:
    return FAISS_INDEX_FILE.exists() and FAISS_METADATA_FILE.exists()


# ── Session state init ────────────────────────────────────────────────────────

def _init_state():
    if "uploaded_files_meta" not in st.session_state:
        st.session_state.uploaded_files_meta = {}
    if "pipeline_status" not in st.session_state:
        st.session_state.pipeline_status = None
    if "pipeline_thread" not in st.session_state:
        st.session_state.pipeline_thread = None
    if "messages" not in st.session_state:
        st.session_state.messages = [
            {"role": "assistant", "content":
             "👋 Bonjour ! Je suis votre assistant documentaire propulsé par **Mistral**.\n\n"
             "Commencez par uploader vos documents dans l'onglet **📂 Sources**, "
             "lancez le pipeline dans **⚙️ Pipeline**, puis revenez ici pour discuter !"}
        ]
    if "rag_model" not in st.session_state:
        st.session_state.rag_model = DEFAULT_MODEL
    if "dark_mode" not in st.session_state:
        st.session_state.dark_mode = True
    if "file_uploader_key" not in st.session_state:
        st.session_state.file_uploader_key = 0
    if "files_validated" not in st.session_state:
        st.session_state.files_validated = False
    if "edit_mode" not in st.session_state:
        st.session_state.edit_mode = False
    if "show_success_popup" not in st.session_state:
        st.session_state.show_success_popup = False

_init_state()


# ── Sidebar ───────────────────────────────────────────────────────────────────

with st.sidebar:
    dark = st.session_state.dark_mode
    _inject_theme_css(dark)

    st.markdown("""
    <div style="text-align:center; padding: 1rem 0 .5rem;">
        <div style="font-size:2.8rem">🧠</div>
        <div style="font-size:1.15rem; font-weight:800; 
             background:linear-gradient(90deg,#6c63ff,#00d4aa);
             -webkit-background-clip:text;-webkit-text-fill-color:transparent;">
            RAG Web App
        </div>
        <div style="color:#8b92a5; font-size:.8rem; margin-top:.2rem;">
            Mistral · FAISS · Docling
        </div>
    </div>
    """, unsafe_allow_html=True)

    theme_label = "☀️ Mode Clair" if dark else "🌙 Mode Sombre"
    if st.button(theme_label, use_container_width=True):
        st.session_state.dark_mode = not dark
        st.rerun()

    st.divider()

    # Statut de l'index
    if index_ready():
        st.markdown('<span class="badge badge-green">✅ Index FAISS prêt</span>', unsafe_allow_html=True)
    else:
        st.markdown('<span class="badge badge-yellow">⚠️ Index non construit</span>', unsafe_allow_html=True)

    n_files = len(st.session_state.uploaded_files_meta)
    st.caption(f"📂 {n_files} fichier(s) chargé(s)")

    st.divider()
    st.markdown("**⚙️ Modèle LLM**")
    model_labels = list(AVAILABLE_MODELS.keys())
    sel_label = st.selectbox(
        "Modèle",
        model_labels,
        index=model_labels.index(
            next((k for k, v in AVAILABLE_MODELS.items() if v == st.session_state.rag_model), model_labels[0])
        ),
        label_visibility="collapsed",
    )
    st.session_state.rag_model = AVAILABLE_MODELS[sel_label]
    st.caption(f"📡 `{st.session_state.rag_model}`")

    st.divider()
    if st.button("🗑️ Effacer la conversation", use_container_width=True):
        st.session_state.messages = [
            {"role": "assistant", "content":
             "👋 Conversation effacée. Comment puis-je vous aider ?"}
        ]
        st.rerun()

    if st.button("🔄 Réinitialiser tout", use_container_width=True, type="secondary"):
        st.session_state.uploaded_files_meta = {}
        st.session_state.pipeline_status = None
        st.session_state.pipeline_thread = None
        st.session_state.files_validated = False
        st.session_state.edit_mode = False
        st.session_state._popup_dismissed = False
        st.rerun()

    st.divider()
    st.caption("🔑 API Mistral : " + ("✅ configurée" if MISTRAL_API_KEY else "❌ manquante"))


# ══════════════════════════════════════════════════════════════════════════════
# ONGLETS
# ══════════════════════════════════════════════════════════════════════════════

tab_sources, tab_pipeline, tab_chat, tab_viz = st.tabs([
    "📂  Sources",
    "⚙️  Pipeline",
    "💬  Chat",
    "📊  Visualisation",
])


# ══════════════════════════════════════════════════════════════════════════════
# ONGLET 1 — SOURCES
# ══════════════════════════════════════════════════════════════════════════════

with tab_sources:

    # ── Hero Upload Zone ─────────────────────────────────────────────────────
    type_tags = " ".join(
        f'<span class="type-tag">{ext}</span>'
        for ext in ["PDF", "DOCX", "PPTX", "XLSX", "HTML", "CSV",
                     "TEX", "PNG", "JPG", "TIFF", "WAV", "MP3", "WebM"]
    )

    st.markdown(f"""
    <div class="hero-upload">
      <div class="hero-icon">📂</div>
      <h1>Gestion des Sources</h1>
      <p class="hero-subtitle">
        Glissez-déposez vos documents ci-dessous ou cliquez pour parcourir vos fichiers
      </p>
      <div class="supported-types">{type_tags}</div>
    </div>
    """, unsafe_allow_html=True)

    # ── Zone d'upload Streamlit (visible if not validated or in edit mode) ──
    is_validated = st.session_state.files_validated
    is_edit_mode = st.session_state.edit_mode

    accepted_types = [
        "pdf", "docx", "pptx", "xlsx",
        "html", "htm", "csv", "tex", "vtt",
        "png", "jpg", "jpeg", "tiff", "tif", "bmp", "webp",
        "wav", "mp3", "m4a", "ogg", "flac", "webm",
    ]

    if not is_validated or is_edit_mode:
        uploaded = st.file_uploader(
            "Glissez-déposez vos fichiers ici",
            type=accepted_types,
            accept_multiple_files=True,
            key=f"file_uploader_{st.session_state.file_uploader_key}",
            help="PDF, DOCX, PPTX, XLSX, HTML, CSV, LaTeX, images (PNG/TIFF/JPEG…), audio (WAV/MP3), WebVTT…",
            label_visibility="collapsed",
        )

        if uploaded:
            INPUTS_DIR.mkdir(parents=True, exist_ok=True)
            new_count = 0
            for uf in uploaded:
                dest = INPUTS_DIR / uf.name
                if uf.name not in st.session_state.uploaded_files_meta:
                    dest.write_bytes(uf.getbuffer())
                    st.session_state.uploaded_files_meta[uf.name] = {
                        "path": dest,
                        "size": uf.size,
                        "suffix": Path(uf.name).suffix.lower(),
                    }
                    new_count += 1
            if new_count:
                st.success(f"✅ {new_count} nouveau(x) fichier(s) ajouté(s)")
                st.rerun()
    else:
        # Show validated status banner
        st.markdown("""
        <div style="text-align:center; padding: 1rem; margin-bottom:1rem;">
            <span class="validated-badge">✅ Fichiers validés — Prêts pour le pipeline</span>
        </div>
        """, unsafe_allow_html=True)

    st.markdown('<hr class="custom-divider">', unsafe_allow_html=True)

    # ── Liste des fichiers ───────────────────────────────────────────────────
    meta = st.session_state.uploaded_files_meta

    if not meta:
        st.session_state.files_validated = False
        st.session_state.edit_mode = False
        st.markdown("""
        <div class="empty-state">
            <div class="empty-icon">📭</div>
            <div class="empty-text">
                Aucun fichier uploadé pour le moment.<br>
                <span style="font-size:0.85rem;">Utilisez la zone ci-dessus pour ajouter vos documents.</span>
            </div>
        </div>
        """, unsafe_allow_html=True)
    else:
        # Stats pills
        total_size = sum(m["size"] for m in meta.values())
        n_docs = sum(1 for m in meta.values() if m["suffix"] in DOCLING_EXTENSIONS)
        n_audio = sum(1 for m in meta.values() if m["suffix"] in AUDIO_EXTENSIONS)

        stats_html = f"""
        <div style="margin-bottom: 1.2rem;">
            <span class="stat-pill">📂 {len(meta)} fichier(s)</span>
            <span class="stat-pill">💾 {human_size(total_size)}</span>
        """
        if n_docs:
            stats_html += f'<span class="stat-pill">📄 {n_docs} document(s)</span>'
        if n_audio:
            stats_html += f'<span class="stat-pill">🎵 {n_audio} audio(s)</span>'
        stats_html += "</div>"
        st.markdown(stats_html, unsafe_allow_html=True)

        col_list, col_preview = st.columns([1, 1], gap="large")

        with col_list:
            header_extra = ' <span class="validated-badge">✅ Validés</span>' if is_validated and not is_edit_mode else ""
            st.markdown(f"""
            <div class="section-header">
                <h3>📋 Fichiers chargés</h3>
                <span class="section-count">{len(meta)}</span>
                {header_extra}
            </div>
            """, unsafe_allow_html=True)

            to_delete = []
            selected_file = st.session_state.get("preview_file")

            for fname, fmeta in meta.items():
                icon = file_icon(fmeta["suffix"])
                size_str = human_size(fmeta["size"])
                is_selected = fname == selected_file
                ext_upper = fmeta["suffix"].replace(".", "").upper()
                validated_class = " validated" if is_validated and not is_edit_mode else ""

                st.markdown(
                    f"""<div class="file-item{validated_class}" style="{'border-color:var(--accent); box-shadow: 0 0 8px rgba(108,99,255,0.15);' if is_selected else ''}">
                    <span class="file-icon">{icon}</span>
                    <span class="file-name">{fname}</span>
                    <span class="file-type-tag">{ext_upper}</span>
                    <span class="file-size">{size_str}</span>
                    </div>""",
                    unsafe_allow_html=True,
                )

                # Show action buttons only if not validated or in edit mode
                if not is_validated or is_edit_mode:
                    btn_cols = st.columns([1, 1])
                    with btn_cols[0]:
                        if st.button("👁️ Aperçu", key=f"prev_{fname}", use_container_width=True):
                            st.session_state.preview_file = fname
                            st.rerun()
                    with btn_cols[1]:
                        if st.button("🗑️ Supprimer", key=f"del_{fname}", use_container_width=True, type="secondary"):
                            to_delete.append(fname)
                else:
                    # Only preview button when validated
                    if st.button("👁️ Aperçu", key=f"prev_{fname}", use_container_width=True):
                        st.session_state.preview_file = fname
                        st.rerun()

            if to_delete:
                for fname in to_delete:
                    fmeta = meta.pop(fname, None)
                    if fmeta and fmeta["path"].exists():
                        fmeta["path"].unlink()
                    if st.session_state.get("preview_file") == fname:
                        st.session_state.preview_file = None
                st.session_state.file_uploader_key += 1
                st.rerun()

            st.markdown('<hr class="custom-divider">', unsafe_allow_html=True)

            # ── Validation / Edit workflow buttons ───────────────────────
            if not is_validated or is_edit_mode:
                # Not yet validated: show Validate + Delete all
                val_col, del_col = st.columns([1, 1])
                with val_col:
                    if st.button("✅ Valider les fichiers", use_container_width=True, type="primary"):
                        st.session_state.files_validated = True
                        st.session_state.edit_mode = False
                        st.rerun()
                with del_col:
                    if st.button("🗑️ Tout supprimer", use_container_width=True, type="secondary"):
                        for fmeta in meta.values():
                            if fmeta["path"].exists():
                                fmeta["path"].unlink()
                        st.session_state.uploaded_files_meta = {}
                        st.session_state.preview_file = None
                        st.session_state.file_uploader_key += 1
                        st.session_state.files_validated = False
                        st.session_state.edit_mode = False
                        st.rerun()
            else:
                # Validated: show Modifier + Go to Pipeline
                mod_col, go_col = st.columns([1, 1])
                with mod_col:
                    if st.button("✏️ Modifier les fichiers", use_container_width=True, type="secondary"):
                        st.session_state.edit_mode = True
                        st.rerun()
                with go_col:
                    st.markdown('<div class="pulse-btn">', unsafe_allow_html=True)
                    if st.button("🚀 Aller au Pipeline →", use_container_width=True, type="primary"):
                        # Navigate to Pipeline tab via JavaScript
                        st.markdown("""
                        <script>
                        const tabs = window.parent.document.querySelectorAll('button[data-baseweb="tab"]');
                        if (tabs.length >= 2) tabs[1].click();
                        </script>
                        """, unsafe_allow_html=True)
                    st.markdown('</div>', unsafe_allow_html=True)

        with col_preview:
            st.markdown("""
            <div class="section-header">
                <h3>🔍 Prévisualisation</h3>
            </div>
            """, unsafe_allow_html=True)

            pfile = st.session_state.get("preview_file")

            if pfile and pfile in meta:
                fmeta = meta[pfile]
                suffix = fmeta["suffix"]
                fpath: Path = fmeta["path"]

                st.markdown(
                    f'<div class="rag-card">'
                    f'<strong>{pfile}</strong> — {human_size(fmeta["size"])}'
                    f' <span class="file-type-tag">{suffix.replace(".", "").upper()}</span>'
                    f'</div>',
                    unsafe_allow_html=True,
                )

                # ── Prévisualisation selon le type ─────────────────────────
                if suffix in {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff", ".tif"}:
                    st.image(str(fpath), use_container_width=True)

                elif suffix == ".pdf":
                    b64 = base64.b64encode(fpath.read_bytes()).decode()
                    st.markdown(
                        f'<div class="preview-box"><iframe src="data:application/pdf;base64,{b64}" '
                        f'width="100%" height="520px" style="border:none;"></iframe></div>',
                        unsafe_allow_html=True,
                    )

                elif suffix in {".html", ".htm"}:
                    content = fpath.read_text(encoding="utf-8", errors="replace")
                    st.components.v1.html(content, height=400, scrolling=True)

                elif suffix in {".wav", ".mp3", ".m4a", ".ogg", ".flac"}:
                    st.audio(str(fpath))

                elif suffix in {".csv"}:
                    try:
                        import pandas as pd
                        df = pd.read_csv(fpath, nrows=100)
                        st.dataframe(df, use_container_width=True, height=350)
                    except Exception as e:
                        st.warning(f"Impossible de lire le CSV : {e}")

                elif suffix in {".xlsx"}:
                    try:
                        import pandas as pd
                        df = pd.read_excel(fpath, nrows=100)
                        st.dataframe(df, use_container_width=True, height=350)
                    except Exception as e:
                        st.warning(f"Impossible de lire le fichier Excel : {e}")

                elif suffix in {".docx"}:
                    try:
                        from docx import Document
                        doc = Document(fpath)
                        text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
                        st.text_area("Contenu DOCX", text[:3000], height=350, disabled=True)
                    except Exception:
                        st.info("⚠️ Prévisualisation DOCX non disponible.")

                elif suffix in {".pptx"}:
                    try:
                        from pptx import Presentation
                        prs = Presentation(fpath)
                        slides_text = []
                        for i, slide in enumerate(prs.slides[:10]):
                            parts = [f"── Slide {i+1} ──"]
                            for shape in slide.shapes:
                                if hasattr(shape, "text") and shape.text.strip():
                                    parts.append(shape.text.strip())
                            slides_text.append("\n".join(parts))
                        st.text_area("Aperçu PPTX (10 premières diapositives)",
                                     "\n\n".join(slides_text), height=350, disabled=True)
                    except Exception:
                        st.info("⚠️ Prévisualisation PPTX non disponible.")

                elif suffix in {".tex", ".vtt", ".webm"}:
                    text = fpath.read_text(encoding="utf-8", errors="replace")
                    st.text_area("Contenu texte", text[:3000], height=350, disabled=True)

                else:
                    st.info(f"Prévisualisation non disponible pour `{suffix}`.")

            else:
                st.markdown("""
                <div class="empty-state">
                    <div class="empty-icon">👁️</div>
                    <div class="empty-text">
                        Cliquez sur <strong>Aperçu</strong> pour prévisualiser un fichier.
                    </div>
                </div>
                """, unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════════════
# ONGLET 2 — PIPELINE
# ══════════════════════════════════════════════════════════════════════════════

with tab_pipeline:
    from pipeline_runner import PipelineStatus, run_pipeline_async

    ps: PipelineStatus | None = st.session_state.pipeline_status
    snap = ps.snapshot() if ps else None

    is_running = bool(snap and snap["running"])
    is_done = bool(snap and snap["done"])
    has_error = bool(snap and snap["error"])

    # ── Pipeline Card (contient tout : titre, étapes, options) ───────────
    st.markdown("""
    <div class="pipeline-card">
      <h1>⚙️ Pipeline RAG</h1>
      <p class="pipeline-subtitle">
        Transformez vos documents en une base de connaissances interrogeable
      </p>
    """, unsafe_allow_html=True)

    # ── Les 4 étapes ─────────────────────────────────────────────────────
    steps_info = [
        ("1", "Extraction",  "Docling + Whisper → Markdown",           "📄"),
        ("2", "Chunking",    "Découpage récursif + chevauchement",      "✂️"),
        ("3", "Embeddings",  "SBERT (+ Mistral optionnel) → vecteurs",  "🧠"),
        ("4", "Indexation",  "FAISS IndexFlatIP → recherche cosine",    "🗄️"),
    ]

    current_step = snap["current_step"] if snap else 0

    # ── Grid layout for steps ────────────────────────────────────────────
    steps_html = '<div class="steps-grid">'
    for num, label, desc, icon in steps_info:
        n = int(num)
        if is_done or (snap and current_step > n):
            css_class = "done"
            indicator = "✅"
        elif snap and current_step == n and is_running:
            css_class = "active"
            indicator = "⏳"
        else:
            css_class = "pending"
            indicator = num

        steps_html += f"""
        <div class="step-card {css_class}">
          <div class="step-icon">{icon}</div>
          <div class="step-num-badge">{indicator}</div>
          <div class="step-card-label">{label}</div>
          <div class="step-card-desc">{desc}</div>
        </div>"""
    steps_html += '</div>'

    st.markdown(steps_html, unsafe_allow_html=True)

    # Options
    st.markdown("</div>", unsafe_allow_html=True)

    # ── Paramètre : Mistral embeddings ───────────────────────────────────
    st.markdown("")

    n_files = len(st.session_state.uploaded_files_meta)

    if n_files == 0:
        st.warning("⚠️ Aucun fichier chargé. Allez d'abord dans **📂 Sources** pour uploader vos documents.")

    use_mistral_emb = st.checkbox(
        "🌐 Utiliser Mistral pour les embeddings (plus puissant, consomme des crédits API)",
        value=False,
        disabled=not MISTRAL_API_KEY,
    )

    # ── Gros bouton Lancer ───────────────────────────────────────────────
    st.markdown("")

    btn_label = "🔄 Relancer le pipeline" if is_done else "🚀 Lancer le pipeline"

    # Status badge
    if is_running:
        st.markdown('<center><span class="badge badge-blue" style="font-size:0.9rem; padding:0.4rem 1rem;">⏳ Pipeline en cours d\'exécution…</span></center>', unsafe_allow_html=True)
    elif is_done:
        st.markdown('<center><span class="badge badge-green" style="font-size:0.9rem; padding:0.4rem 1rem;">✅ Pipeline terminé avec succès</span></center>', unsafe_allow_html=True)
    elif has_error:
        st.markdown('<center><span class="badge badge-red" style="font-size:0.9rem; padding:0.4rem 1rem;">❌ Erreur lors de l\'exécution</span></center>', unsafe_allow_html=True)

    st.markdown("")
    # Use pulse animation on the launch button
    st.markdown('<div class="launch-btn-wrap pulse-btn">', unsafe_allow_html=True)
    if st.button(
        btn_label,
        disabled=(n_files == 0 or is_running),
        type="primary",
        use_container_width=True,
    ):
        new_ps = PipelineStatus()
        st.session_state.pipeline_status = new_ps
        file_paths = [m["path"] for m in st.session_state.uploaded_files_meta.values()]
        t = run_pipeline_async(file_paths, new_ps, use_mistral_embed=use_mistral_emb)
        st.session_state.pipeline_thread = t
        st.session_state.show_success_popup = False
        time.sleep(0.3)
        st.rerun()
    st.markdown("</div>", unsafe_allow_html=True)

    # ── Barre de progression ─────────────────────────────────────────────
    st.markdown("")

    if snap:
        if is_running:
            overall = (snap["current_step"] - 1) / snap["total_steps"] + \
                      (snap["step_progress"] / max(snap["step_total"], 1)) / snap["total_steps"]
            st.progress(min(overall, 0.99), text=f"Étape {snap['current_step']}/{snap['total_steps']} — {snap['step_message']}")

        elif is_done:
            st.progress(1.0, text="✅ Pipeline terminé !")
            st.success(f"🎉 Index construit avec succès — **{snap['total_vectors']} vecteurs** indexés")

            # Résumé post-pipeline
            m_col1, m_col2, m_col3 = st.columns(3)
            with m_col1:
                st.metric("📂 Fichiers traités", len(st.session_state.uploaded_files_meta))
            with m_col2:
                st.metric("🗄️ Vecteurs indexés", snap["total_vectors"])
            with m_col3:
                st.metric("📡 Embedding", "SBERT" + (" + Mistral" if use_mistral_emb else ""))

            # ── Success banner at bottom ──────────────────────────────────
            st.markdown("""
            <div style="margin-top:1.5rem; padding:1.5rem 2rem; border-radius:16px;
                        background: linear-gradient(135deg, rgba(0,169,143,0.12), rgba(108,99,255,0.08));
                        border: 1px solid rgba(0,212,170,0.3); text-align:center;">
                <div style="font-size:2.5rem; margin-bottom:0.5rem;">🎉</div>
                <div style="font-size:1.3rem; font-weight:800;
                            background:linear-gradient(90deg,#00a98f,var(--accent));
                            -webkit-background-clip:text; -webkit-text-fill-color:transparent;
                            margin-bottom:0.4rem;">Pipeline terminé avec succès !</div>
                <div style="color:var(--text-muted); margin-bottom:1rem;">Votre base de connaissances est prête. Discutez avec vos documents dès maintenant !</div>
            </div>
            """, unsafe_allow_html=True)

            st.markdown('<div class="pulse-btn" style="margin-top:1rem;">', unsafe_allow_html=True)
            if st.button("💬 Aller au Chat →", use_container_width=True, type="primary", key="go_chat_btn"):
                # Navigate to Chat tab via JavaScript
                st.markdown("""
                <script>
                const tabs = window.parent.document.querySelectorAll('button[data-baseweb="tab"]');
                if (tabs.length >= 3) tabs[2].click();
                </script>
                """, unsafe_allow_html=True)
            st.markdown('</div>', unsafe_allow_html=True)

        elif has_error:
            st.error(snap["error"])
            st.progress(0.0)
    else:
        st.progress(0.0, text="En attente du lancement…")

    # ── Journal d'exécution (logs en temps réel) ─────────────────────────
    st.markdown("")
    st.markdown("""
    <div class="section-header">
        <h3>🖥️ Journal d'exécution</h3>
    </div>
    """, unsafe_allow_html=True)

    if snap and snap["logs"]:
        log_lines = snap["logs"][-100:]
        log_text = "\n".join(log_lines)
        st.markdown(f'<div class="log-box">{log_text}</div>', unsafe_allow_html=True)
    else:
        st.markdown(
            '<div class="log-box" style="text-align:center; color:var(--text-muted); padding:2rem;">'
            '📋 Les logs apparaîtront ici une fois le pipeline lancé…</div>',
            unsafe_allow_html=True,
        )

    # Auto-refresh pendant l'exécution pour synchroniser les logs en temps réel
    if is_running:
        time.sleep(1.0)
        st.rerun()


# ══════════════════════════════════════════════════════════════════════════════
# ONGLET 3 — CHAT
# ══════════════════════════════════════════════════════════════════════════════

with tab_chat:

    # ── Header Chat ──────────────────────────────────────────────────────
    st.markdown("""
    <div class="chat-welcome">
      <h1>💬 Assistant RAG</h1>
      <p class="chat-subtitle">
        Posez vos questions sur les documents indexés — les réponses citent leurs sources
      </p>
      <div class="chat-feature-grid">
        <div class="chat-feature">
          <div class="feat-icon">🔍</div>
          <div class="feat-title">Recherche sémantique</div>
          <div class="feat-desc">Trouve les passages pertinents dans vos documents</div>
        </div>
        <div class="chat-feature">
          <div class="feat-icon">📚</div>
          <div class="feat-title">Sources citées</div>
          <div class="feat-desc">Chaque réponse cite les documents utilisés</div>
        </div>
        <div class="chat-feature">
          <div class="feat-icon">🧠</div>
          <div class="feat-title">Mistral AI</div>
          <div class="feat-desc">Propulsé par les modèles Mistral de dernière génération</div>
        </div>
      </div>
    </div>
    """, unsafe_allow_html=True)

    # ── Vérification de l'index ──────────────────────────────────────────
    if not index_ready():
        st.info(
            "💡 **L'index FAISS n'est pas encore construit.**\n\n"
            "Pour l'utiliser :\n"
            "1. Uploadez vos documents dans **📂 Sources**\n"
            "2. Lancez le traitement dans **⚙️ Pipeline**\n"
            "3. Revenez ici pour discuter !\n\n"
            "_Le chat fonctionne également sans index (mode conversationnel)._"
        )

    # ── Fonctions LLM + RAG ──────────────────────────────────────────────
    def get_rag_context(question: str, top_k: int = 4) -> tuple[str, list[str]]:
        """Recherche dans l'index FAISS et retourne (contexte_texte, sources)."""
        if not index_ready():
            return "", []
        try:
            from src.step_4_store.vector_store import search
            results = search(question, top_k=top_k)
            if not results:
                return "", []

            texte = ""
            sources = []
            for r in results:
                texte += (
                    f"---\nDocument : {r['metadata']['source']}\n"
                    f"Extrait : {r['text']}\n"
                )
                src = r["metadata"]["source"].split("/")[-1]
                if src not in sources:
                    sources.append(src)
            return texte, sources
        except Exception as e:
            logger.error("Erreur RAG : %s", e)
            return "", []

    def build_messages(history: list[dict], question: str, context: str) -> list:
        """Construit les messages formatés pour l'API Mistral."""
        from mistralai.models import UserMessage, AssistantMessage, SystemMessage

        system_prompt = """\
### RÔLE :
Vous êtes un assistant documentaire intelligent et bienveillant.
Vous répondez uniquement à partir des documents fournis par l'utilisateur.

### COMPORTEMENT :
- Répondez en français, de façon claire, structurée et précise.
- Citez toujours vos sources (noms de fichiers) si disponibles.
- Si l'information n'est pas dans les documents : dites-le clairement, ne l'inventez pas.
- Utilisez des listes et titres Markdown pour les réponses longues.

### INTERDICTIONS :
- Ne jamais inventer de faits.
- Ne jamais répondre sur des sujets non couverts par les documents.
"""

        formatted = [SystemMessage(content=system_prompt)]

        recent = history[-8:] if len(history) > 8 else history
        for msg in recent[:-1]:
            if msg["role"] == "user":
                formatted.append(UserMessage(content=msg["content"]))
            elif msg["role"] == "assistant":
                formatted.append(AssistantMessage(content=msg["content"]))

        if context:
            enriched = (
                f"### DOCUMENTS PERTINENTS :\n{context}\n\n"
                f"### QUESTION :\n{question}"
            )
        else:
            enriched = question

        formatted.append(UserMessage(content=enriched))
        return formatted

    def stream_llm(messages: list, model: str):
        """Appelle l'API Mistral en mode streaming et yield les tokens."""
        if not MISTRAL_API_KEY:
            yield "❌ Clé API Mistral non configurée. Veuillez définir `MISTRAL_API_KEY` dans le fichier `.env`."
            return
        try:
            from mistralai import Mistral
            client = Mistral(api_key=MISTRAL_API_KEY)
            stream_resp = client.chat.stream(
                model=model,
                messages=messages,
                **GENERATION_PARAMS,
            )
            for chunk in stream_resp:
                token = chunk.data.choices[0].delta.content
                if token:
                    yield token
        except Exception as e:
            logger.error("Erreur LLM : %s", e)
            yield f"❌ Erreur lors de l'appel à Mistral : {e}"

    # ── Affichage de l'historique ────────────────────────────────────────
    for msg in st.session_state.messages:
        with st.chat_message(msg["role"]):
            st.markdown(msg["content"])

    # ── Saisie utilisateur ───────────────────────────────────────────────
    if prompt := st.chat_input("Posez votre question sur vos documents…"):
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)

        # Recherche RAG
        with st.spinner("🔍 Recherche dans les documents…"):
            ctx_text, sources = get_rag_context(prompt)

        # Construire les messages
        prompt_messages = build_messages(st.session_state.messages, prompt, ctx_text)

        # Générer la réponse en streaming (mot par mot)
        with st.chat_message("assistant"):
            answer = st.write_stream(stream_llm(prompt_messages, st.session_state.rag_model))

            if sources:
                sources_md = ", ".join([f"`{s}`" for s in sources])
                st.caption(f"📚 **Sources utilisées :** {sources_md}")
                answer += f"\n\n*📚 Sources : {', '.join(sources)}*"
            elif not ctx_text and index_ready():
                st.caption("ℹ️ Aucun document pertinent trouvé pour cette question.")

        st.session_state.messages.append({"role": "assistant", "content": answer})


# ══════════════════════════════════════════════════════════════════════════════
# ONGLET 4 — VISUALISATION
# ══════════════════════════════════════════════════════════════════════════════

with tab_viz:

    st.markdown("""
    <div class="pipeline-card">
      <h1>📊 Visualisation des Résultats</h1>
      <p class="pipeline-subtitle">
        Explorez les embeddings, chunks et l'index vectoriel après l'exécution du pipeline
      </p>
    </div>
    """, unsafe_allow_html=True)

    if not index_ready():
        st.info(
            "💡 **Aucune donnée à visualiser pour le moment.**\n\n"
            "Lancez d'abord le pipeline dans l'onglet **⚙️ Pipeline** pour générer les données."
        )
    else:
        import json as _json

        # ── Load data ────────────────────────────────────────────────────
        viz_chunks_data = None
        viz_chunks = []
        viz_metadata = []
        viz_index = None
        viz_embeddings = None

        chunks_file = ROOT_DIR / "data" / "chunks" / "chunks.json"
        embeddings_file = ROOT_DIR / "data" / "chunks" / "embeddings.npz"

        if chunks_file.exists():
            with open(chunks_file, "r", encoding="utf-8") as f:
                viz_chunks_data = _json.load(f)
                viz_chunks = viz_chunks_data.get("chunks", [])

        if FAISS_METADATA_FILE.exists():
            with open(FAISS_METADATA_FILE, "r", encoding="utf-8") as f:
                viz_metadata = _json.load(f)

        try:
            import faiss as _faiss
            if FAISS_INDEX_FILE.exists():
                viz_index = _faiss.read_index(str(FAISS_INDEX_FILE))
        except Exception:
            pass

        if embeddings_file.exists():
            try:
                import numpy as _np
                data = _np.load(str(embeddings_file))
                if "sbert" in data:
                    viz_embeddings = data["sbert"]
                elif len(data.files) > 0:
                    viz_embeddings = data[data.files[0]]
            except Exception:
                pass

        # ── Stats cards ──────────────────────────────────────────────────
        st.markdown("")
        c1, c2, c3, c4 = st.columns(4)
        with c1:
            st.markdown(f"""
            <div class="viz-card">
                <h3>🗄️ Index FAISS</h3>
                <div class="viz-stat"><span class="viz-stat-label">Vecteurs</span><span class="viz-stat-value">{viz_index.ntotal if viz_index else 0}</span></div>
                <div class="viz-stat"><span class="viz-stat-label">Dimensions</span><span class="viz-stat-value">{viz_index.d if viz_index else '—'}</span></div>
            </div>""", unsafe_allow_html=True)
        with c2:
            st.markdown(f"""
            <div class="viz-card">
                <h3>✂️ Chunks</h3>
                <div class="viz-stat"><span class="viz-stat-label">Total</span><span class="viz-stat-value">{len(viz_chunks)}</span></div>
                <div class="viz-stat"><span class="viz-stat-label">Config</span><span class="viz-stat-value">{viz_chunks_data.get('config', {}).get('chunk_size', '—') if viz_chunks_data else '—'}/{viz_chunks_data.get('config', {}).get('chunk_overlap', '—') if viz_chunks_data else '—'}</span></div>
            </div>""", unsafe_allow_html=True)
        with c3:
            n_sources = len(set(c.get("metadata", {}).get("filename", "") for c in viz_chunks)) if viz_chunks else 0
            st.markdown(f"""
            <div class="viz-card">
                <h3>📂 Sources</h3>
                <div class="viz-stat"><span class="viz-stat-label">Documents</span><span class="viz-stat-value">{n_sources}</span></div>
                <div class="viz-stat"><span class="viz-stat-label">Métadonnées</span><span class="viz-stat-value">{len(viz_metadata)}</span></div>
            </div>""", unsafe_allow_html=True)
        with c4:
            emb_shape = viz_embeddings.shape if viz_embeddings is not None else None
            st.markdown(f"""
            <div class="viz-card">
                <h3>🧠 Embeddings</h3>
                <div class="viz-stat"><span class="viz-stat-label">Shape</span><span class="viz-stat-value">{f'{emb_shape[0]}×{emb_shape[1]}' if emb_shape else '—'}</span></div>
                <div class="viz-stat"><span class="viz-stat-label">Modèle</span><span class="viz-stat-value">SBERT</span></div>
            </div>""", unsafe_allow_html=True)

        st.markdown('<hr class="custom-divider">', unsafe_allow_html=True)

        # ── Embedding Visualization ──────────────────────────────────────
        st.markdown("""
        <div class="section-header">
            <h3>🧠 Visualisation des Embeddings (2D)</h3>
        </div>
        """, unsafe_allow_html=True)

        if viz_embeddings is not None and len(viz_embeddings) >= 2:
            import numpy as _np
            try:
                from sklearn.manifold import TSNE
                from sklearn.decomposition import PCA
                import matplotlib
                matplotlib.use("Agg")
                import matplotlib.pyplot as plt

                n_samples = len(viz_embeddings)
                if n_samples < 5:
                    reducer = PCA(n_components=2, random_state=42)
                    coords = reducer.fit_transform(viz_embeddings)
                    method = "PCA"
                else:
                    perp = min(30, n_samples - 1)
                    reducer = TSNE(n_components=2, perplexity=perp, random_state=42, init="pca", learning_rate="auto")
                    coords = reducer.fit_transform(viz_embeddings)
                    method = "t-SNE"

                filenames = [c.get("metadata", {}).get("filename", "Inconnu") for c in viz_chunks[:n_samples]]

                fig, ax = plt.subplots(figsize=(12, 7))
                fig.patch.set_facecolor('#0d0f14' if st.session_state.dark_mode else '#f5f7fb')
                ax.set_facecolor('#161b27' if st.session_state.dark_mode else '#ffffff')

                unique_files = list(set(filenames))
                colors = plt.cm.Set2.colors if len(unique_files) <= 8 else plt.cm.tab20.colors
                color_map = {f: colors[i % len(colors)] for i, f in enumerate(unique_files)}

                for fname in unique_files:
                    mask = [i for i, fn in enumerate(filenames) if fn == fname]
                    ax.scatter(
                        coords[mask, 0], coords[mask, 1],
                        label=fname, alpha=0.75, s=80, edgecolors='white', linewidths=0.5,
                        color=color_map[fname],
                    )
                text_color = '#e8eaf0' if st.session_state.dark_mode else '#1a1d2e'
                ax.set_title(f"Embeddings ({method}) — {n_samples} chunks", fontsize=14, fontweight="bold", color=text_color, pad=15)
                ax.tick_params(colors=text_color)
                for spine in ax.spines.values():
                    spine.set_color('#2d3148' if st.session_state.dark_mode else '#d0d5e8')
                ax.legend(fontsize=8, loc="upper right", framealpha=0.7)
                fig.tight_layout()
                st.pyplot(fig)
                plt.close(fig)
            except Exception as e:
                st.warning(f"Impossible de générer la visualisation : {e}")
        else:
            st.info("Pas assez de données pour la visualisation 2D.")

        st.markdown('<hr class="custom-divider">', unsafe_allow_html=True)

        # ── Chunks Explorer ──────────────────────────────────────────────
        st.markdown("""
        <div class="section-header">
            <h3>✂️ Explorateur de Chunks</h3>
        </div>
        """, unsafe_allow_html=True)

        if viz_chunks:
            # Source filter
            all_sources = sorted(set(c.get("metadata", {}).get("filename", "Inconnu") for c in viz_chunks))
            selected_source = st.selectbox("🔍 Filtrer par source", ["Tous"] + all_sources, key="viz_source_filter")

            filtered = viz_chunks if selected_source == "Tous" else [
                c for c in viz_chunks if c.get("metadata", {}).get("filename") == selected_source
            ]

            st.caption(f"Affichage de {min(20, len(filtered))}/{len(filtered)} chunks")

            for i, chunk in enumerate(filtered[:20]):
                meta = chunk.get("metadata", {})
                text_preview = chunk.get("text", "")[:200].replace("\n", " ")
                st.markdown(f"""
                <div class="viz-card">
                    <div style="display:flex; justify-content:space-between; align-items:center; margin-bottom:0.5rem;">
                        <span style="font-weight:700; font-size:0.9rem;">Chunk #{meta.get('chunk_index', i)}</span>
                        <span class="file-type-tag">{meta.get('filename', 'Inconnu')}</span>
                    </div>
                    <div style="color:var(--text-muted); font-size:0.85rem; line-height:1.5;">{text_preview}…</div>
                    <div style="margin-top:0.5rem; display:flex; gap:1rem;">
                        <span class="stat-pill">📏 {meta.get('chunk_size', len(chunk.get('text', '')))} car.</span>
                        <span class="stat-pill">📂 {meta.get('category', '—')}</span>
                    </div>
                </div>
                """, unsafe_allow_html=True)
        else:
            st.info("Aucun chunk disponible. Lancez le pipeline d'abord.")

