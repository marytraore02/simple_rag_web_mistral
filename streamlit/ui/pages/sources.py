"""
Page Sources — Upload, prévisualisation et validation des fichiers.
"""

from __future__ import annotations

import base64
from pathlib import Path

import streamlit as st

from app_config.settings import INPUTS_DIR, DOCLING_EXTENSIONS, AUDIO_EXTENSIONS
from utils.helpers import human_size, file_icon


def render(tab) -> None:
    """Affiche l'onglet Sources dans le conteneur `tab`."""

    with tab:
        # ── Hero Upload Zone ─────────────────────────────────────────────
        type_tags = " ".join(
            f'<span class="type-tag">{ext}</span>'
            for ext in ["PDF", "DOCX", "PPTX", "XLSX", "HTML", "CSV",
                         "TEX", "PNG", "JPG", "TIFF", "WAV", "MP3", "WebM"]
        )

        st.markdown(f"""
        <div class="hero-upload">
          <div class="hero-icon">📂</div>
          <h1>Gestion des Sources</h1>
          <p class="hero-subtitle">
            Glissez-déposez vos documents ci-dessous ou cliquez pour parcourir vos fichiers
          </p>
          <div class="supported-types">{type_tags}</div>
        </div>
        """, unsafe_allow_html=True)

        # ── Zone d'upload (visible si non validé ou en mode édition) ─────
        is_validated = st.session_state.files_validated
        is_edit_mode = st.session_state.edit_mode

        accepted_types = [
            "pdf", "docx", "pptx", "xlsx",
            "html", "htm", "csv", "tex", "vtt",
            "png", "jpg", "jpeg", "tiff", "tif", "bmp", "webp",
            "wav", "mp3", "m4a", "ogg", "flac", "webm",
        ]

        if not is_validated or is_edit_mode:
            uploaded = st.file_uploader(
                "Glissez-déposez vos fichiers ici",
                type=accepted_types,
                accept_multiple_files=True,
                key=f"file_uploader_{st.session_state.file_uploader_key}",
                help="PDF, DOCX, PPTX, XLSX, HTML, CSV, LaTeX, images, audio, WebVTT…",
                label_visibility="collapsed",
            )

            if uploaded:
                INPUTS_DIR.mkdir(parents=True, exist_ok=True)
                new_count = 0
                for uf in uploaded:
                    dest = INPUTS_DIR / uf.name
                    if uf.name not in st.session_state.uploaded_files_meta:
                        dest.write_bytes(uf.getbuffer())
                        st.session_state.uploaded_files_meta[uf.name] = {
                            "path": dest,
                            "size": uf.size,
                            "suffix": Path(uf.name).suffix.lower(),
                        }
                        new_count += 1
                if new_count:
                    st.success(f"✅ {new_count} nouveau(x) fichier(s) ajouté(s)")
                    st.rerun()
        else:
            st.markdown("""
            <div style="text-align:center; padding: 1rem; margin-bottom:1rem;">
                <span class="validated-badge">✅ Fichiers validés — Prêts pour le pipeline</span>
            </div>
            """, unsafe_allow_html=True)

        st.markdown('<hr class="custom-divider">', unsafe_allow_html=True)

        # ── Liste des fichiers ───────────────────────────────────────────
        meta = st.session_state.uploaded_files_meta

        if not meta:
            st.session_state.files_validated = False
            st.session_state.edit_mode = False
            st.markdown("""
            <div class="empty-state">
                <div class="empty-icon">📭</div>
                <div class="empty-text">
                    Aucun fichier uploadé pour le moment.<br>
                    <span style="font-size:0.85rem;">Utilisez la zone ci-dessus pour ajouter vos documents.</span>
                </div>
            </div>
            """, unsafe_allow_html=True)
        else:
            _render_file_list(meta, is_validated, is_edit_mode)


def _render_file_list(meta: dict, is_validated: bool, is_edit_mode: bool) -> None:
    """Affiche la liste des fichiers et le panneau de prévisualisation."""

    # Stats pills
    total_size = sum(m["size"] for m in meta.values())
    n_docs = sum(1 for m in meta.values() if m["suffix"] in DOCLING_EXTENSIONS)
    n_audio = sum(1 for m in meta.values() if m["suffix"] in AUDIO_EXTENSIONS)

    stats_html = f"""
    <div style="margin-bottom: 1.2rem;">
        <span class="stat-pill">📂 {len(meta)} fichier(s)</span>
        <span class="stat-pill">💾 {human_size(total_size)}</span>
    """
    if n_docs:
        stats_html += f'<span class="stat-pill">📄 {n_docs} document(s)</span>'
    if n_audio:
        stats_html += f'<span class="stat-pill">🎵 {n_audio} audio(s)</span>'
    stats_html += "</div>"
    st.markdown(stats_html, unsafe_allow_html=True)

    col_list, col_preview = st.columns([1, 1], gap="large")

    with col_list:
        _render_file_column(meta, is_validated, is_edit_mode)

    with col_preview:
        _render_preview_column(meta)


def _render_file_column(meta: dict, is_validated: bool, is_edit_mode: bool) -> None:
    """Colonne gauche : liste des fichiers + boutons d'actions."""

    header_extra = ' <span class="validated-badge">✅ Validés</span>' if is_validated and not is_edit_mode else ""
    st.markdown(f"""
    <div class="section-header">
        <h3>📋 Fichiers chargés</h3>
        <span class="section-count">{len(meta)}</span>
        {header_extra}
    </div>
    """, unsafe_allow_html=True)

    to_delete = []
    selected_file = st.session_state.get("preview_file")

    for fname, fmeta in meta.items():
        icon = file_icon(fmeta["suffix"])
        size_str = human_size(fmeta["size"])
        is_selected = fname == selected_file
        ext_upper = fmeta["suffix"].replace(".", "").upper()
        validated_class = " validated" if is_validated and not is_edit_mode else ""

        st.markdown(
            f"""<div class="file-item{validated_class}" style="{'border-color:var(--accent); box-shadow: 0 0 8px rgba(108,99,255,0.15);' if is_selected else ''}">
            <span class="file-icon">{icon}</span>
            <span class="file-name">{fname}</span>
            <span class="file-type-tag">{ext_upper}</span>
            <span class="file-size">{size_str}</span>
            </div>""",
            unsafe_allow_html=True,
        )

        if not is_validated or is_edit_mode:
            btn_cols = st.columns([1, 1])
            with btn_cols[0]:
                if st.button("👁️ Aperçu", key=f"prev_{fname}", use_container_width=True):
                    st.session_state.preview_file = fname
                    st.rerun()
            with btn_cols[1]:
                if st.button("🗑️ Supprimer", key=f"del_{fname}", use_container_width=True, type="secondary"):
                    to_delete.append(fname)
        else:
            if st.button("👁️ Aperçu", key=f"prev_{fname}", use_container_width=True):
                st.session_state.preview_file = fname
                st.rerun()

    # Handle deletions
    if to_delete:
        for fname in to_delete:
            fmeta = meta.pop(fname, None)
            if fmeta and fmeta["path"].exists():
                fmeta["path"].unlink()
            if st.session_state.get("preview_file") == fname:
                st.session_state.preview_file = None
        st.session_state.file_uploader_key += 1
        st.rerun()

    st.markdown('<hr class="custom-divider">', unsafe_allow_html=True)

    # ── Validation / Edit workflow buttons ────────────────────────────
    if not is_validated or is_edit_mode:
        val_col, del_col = st.columns([1, 1])
        with val_col:
            if st.button("✅ Valider les fichiers", use_container_width=True, type="primary"):
                st.session_state.files_validated = True
                st.session_state.edit_mode = False
                st.rerun()
        with del_col:
            if st.button("🗑️ Tout supprimer", use_container_width=True, type="secondary"):
                for fmeta in meta.values():
                    if fmeta["path"].exists():
                        fmeta["path"].unlink()
                st.session_state.uploaded_files_meta = {}
                st.session_state.preview_file = None
                st.session_state.file_uploader_key += 1
                st.session_state.files_validated = False
                st.session_state.edit_mode = False
                st.rerun()
    else:
        mod_col, go_col = st.columns([1, 1])
        with mod_col:
            if st.button("✏️ Modifier les fichiers", use_container_width=True, type="secondary"):
                st.session_state.edit_mode = True
                st.rerun()
        with go_col:
            st.markdown('<div class="pulse-btn">', unsafe_allow_html=True)
            if st.button("🚀 Aller au Pipeline →", use_container_width=True, type="primary"):
                st.components.v1.html("""
                <script>
                const tabs = window.parent.document.querySelectorAll('[data-baseweb="tab"]');
                if (tabs.length >= 2) { tabs[1].click(); }
                </script>
                """, height=0)
            st.markdown('</div>', unsafe_allow_html=True)


def _render_preview_column(meta: dict) -> None:
    """Colonne droite : prévisualisation du fichier sélectionné."""

    st.markdown("""
    <div class="section-header">
        <h3>🔍 Prévisualisation</h3>
    </div>
    """, unsafe_allow_html=True)

    pfile = st.session_state.get("preview_file")

    if pfile and pfile in meta:
        fmeta = meta[pfile]
        suffix = fmeta["suffix"]
        fpath: Path = fmeta["path"]

        st.markdown(
            f'<div class="rag-card">'
            f'<strong>{pfile}</strong> — {human_size(fmeta["size"])}'
            f' <span class="file-type-tag">{suffix.replace(".", "").upper()}</span>'
            f'</div>',
            unsafe_allow_html=True,
        )

        # Prévisualisation selon le type
        if suffix in {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tiff", ".tif"}:
            st.image(str(fpath), use_container_width=True)
        elif suffix == ".pdf":
            b64 = base64.b64encode(fpath.read_bytes()).decode()
            st.markdown(
                f'<div class="preview-box"><iframe src="data:application/pdf;base64,{b64}" '
                f'width="100%" height="520px" style="border:none;"></iframe></div>',
                unsafe_allow_html=True,
            )
        elif suffix in {".html", ".htm"}:
            content = fpath.read_text(encoding="utf-8", errors="replace")
            st.components.v1.html(content, height=400, scrolling=True)
        elif suffix in {".wav", ".mp3", ".m4a", ".ogg", ".flac"}:
            st.audio(str(fpath))
        elif suffix == ".csv":
            try:
                import pandas as pd
                df = pd.read_csv(fpath, nrows=100)
                st.dataframe(df, use_container_width=True, height=350)
            except Exception as e:
                st.warning(f"Impossible de lire le CSV : {e}")
        elif suffix == ".xlsx":
            try:
                import pandas as pd
                df = pd.read_excel(fpath, nrows=100)
                st.dataframe(df, use_container_width=True, height=350)
            except Exception as e:
                st.warning(f"Impossible de lire le fichier Excel : {e}")
        elif suffix == ".docx":
            try:
                from docx import Document
                doc = Document(fpath)
                text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
                st.text_area("Contenu DOCX", text[:3000], height=350, disabled=True)
            except Exception:
                st.info("⚠️ Prévisualisation DOCX non disponible.")
        elif suffix == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(fpath)
                slides_text = []
                for i, slide in enumerate(prs.slides[:10]):
                    parts = [f"── Slide {i+1} ──"]
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            parts.append(shape.text.strip())
                    slides_text.append("\n".join(parts))
                st.text_area("Aperçu PPTX (10 premières diapositives)",
                             "\n\n".join(slides_text), height=350, disabled=True)
            except Exception:
                st.info("⚠️ Prévisualisation PPTX non disponible.")
        elif suffix in {".tex", ".vtt", ".webm"}:
            text = fpath.read_text(encoding="utf-8", errors="replace")
            st.text_area("Contenu texte", text[:3000], height=350, disabled=True)
        else:
            st.info(f"Prévisualisation non disponible pour `{suffix}`.")
    else:
        st.markdown("""
        <div class="empty-state">
            <div class="empty-icon">👁️</div>
            <div class="empty-text">
                Cliquez sur <strong>Aperçu</strong> pour prévisualiser un fichier.
            </div>
        </div>
        """, unsafe_allow_html=True)
